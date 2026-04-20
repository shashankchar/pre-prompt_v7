from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = r"C:\Users\LENOVO\Desktop\prompt bank\PrePrompt_5_Slides_UseCase.pptx"


BG_DARK = RGBColor(9, 10, 14)
ACCENT_RED = RGBColor(229, 9, 20)
ACCENT_ORANGE = RGBColor(255, 123, 0)
TEXT_LIGHT = RGBColor(245, 245, 245)
TEXT_MUTED = RGBColor(200, 200, 200)


def style_slide_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def style_title(shape, size=40):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = TEXT_LIGHT


def style_body_text(shape, size=22):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = TEXT_MUTED


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    style_slide_background(slide)
    slide.shapes.title.text = "PrePrompt"
    slide.placeholders[1].text = "Prompt Bank + PromptTyper\nPresentation (5 Slides)"
    style_title(slide.shapes.title, size=50)
    style_body_text(slide.placeholders[1], size=24)


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    style_slide_background(slide)
    slide.shapes.title.text = title
    style_title(slide.shapes.title, size=40)

    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_MUTED


def add_system_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    style_slide_background(slide)
    slide.shapes.title.text = "System Workflow"
    style_title(slide.shapes.title, size=40)

    top = Inches(2.2)
    box_w = Inches(2.2)
    box_h = Inches(1.0)
    gap = Inches(0.35)
    start_left = Inches(0.5)

    labels = [
        "User",
        "Prompt Bank\n(Web UI)",
        "Bridge API\n(127.0.0.1:8765)",
        "PromptTyper\n(Manager)",
        "Espanso\nTyping Engine",
    ]

    boxes = []
    for i, label in enumerate(labels):
        left = start_left + i * (box_w + gap)
        rect = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, box_w, box_h
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(24, 24, 30)
        rect.line.color.rgb = ACCENT_RED if i in (1, 3) else RGBColor(90, 90, 98)
        rect.line.width = Pt(1.5)
        rect.text_frame.text = label
        rect.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in rect.text_frame.paragraphs[0].runs:
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = TEXT_LIGHT
        boxes.append(rect)

    for i in range(len(boxes) - 1):
        x1 = boxes[i].left + boxes[i].width
        y = boxes[i].top + boxes[i].height // 2
        x2 = boxes[i + 1].left
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
        line.line.color.rgb = ACCENT_ORANGE
        line.line.width = Pt(2)


def add_use_case_diagram_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    style_slide_background(slide)
    slide.shapes.title.text = "Use Case Diagram"
    style_title(slide.shapes.title, size=40)

    # System boundary
    boundary = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.3), Inches(1.5), Inches(8.4), Inches(4.8)
    )
    boundary.fill.solid()
    boundary.fill.fore_color.rgb = RGBColor(18, 18, 24)
    boundary.line.color.rgb = RGBColor(110, 110, 120)
    boundary.text_frame.text = "PrePrompt System"
    boundary.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in boundary.text_frame.paragraphs[0].runs:
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = TEXT_MUTED

    # Actors
    user = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2), Inches(1.5), Inches(0.9)
    )
    user.fill.solid()
    user.fill.fore_color.rgb = RGBColor(35, 35, 43)
    user.line.color.rgb = ACCENT_ORANGE
    user.text_frame.text = "User"
    user.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in user.text_frame.paragraphs[0].runs:
        run.font.bold = True
        run.font.color.rgb = TEXT_LIGHT

    prompttyper = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.0), Inches(2.2), Inches(1.8), Inches(0.9)
    )
    prompttyper.fill.solid()
    prompttyper.fill.fore_color.rgb = RGBColor(35, 35, 43)
    prompttyper.line.color.rgb = ACCENT_RED
    prompttyper.text_frame.text = "PromptTyper"
    prompttyper.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in prompttyper.text_frame.paragraphs[0].runs:
        run.font.bold = True
        run.font.color.rgb = TEXT_LIGHT

    # Use cases
    use_cases = [
        ("Browse/Search Prompts", Inches(3.0), Inches(2.0)),
        ("Test Prompt", Inches(5.2), Inches(2.0)),
        ("Save to PromptTyper", Inches(7.4), Inches(2.0)),
        ("Manage Prompts", Inches(4.1), Inches(3.6)),
        ("Export/Import JSON", Inches(6.5), Inches(3.6)),
    ]

    ellipses = []
    for text, left, top in use_cases:
        e = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, left, top, Inches(2.2), Inches(1.0)
        )
        e.fill.solid()
        e.fill.fore_color.rgb = RGBColor(30, 30, 38)
        e.line.color.rgb = RGBColor(170, 170, 180)
        e.text_frame.text = text
        e.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in e.text_frame.paragraphs[0].runs:
            run.font.size = Pt(12)
            run.font.color.rgb = TEXT_LIGHT
        ellipses.append(e)

    # Connectors from actors
    for idx in (0, 1, 3):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            user.left + user.width,
            user.top + user.height // 2,
            ellipses[idx].left,
            ellipses[idx].top + ellipses[idx].height // 2,
        )
        line.line.color.rgb = ACCENT_ORANGE
        line.line.width = Pt(1.8)

    for idx in (2, 3):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            ellipses[idx].left + ellipses[idx].width,
            ellipses[idx].top + ellipses[idx].height // 2,
            prompttyper.left,
            prompttyper.top + prompttyper.height // 2,
        )
        line.line.color.rgb = ACCENT_RED
        line.line.width = Pt(1.8)


def add_results_slide(prs):
    add_bullet_slide(prs, "Results and Conclusion", [
        "Prompt creation, testing, and saving are now in one smooth flow.",
        "Non-technical users can manage prompts without touching code.",
        "PromptTyper integration reduces manual copy-paste effort.",
        "Project is ready for demo and practical daily use.",
    ])


def build_presentation():
    prs = Presentation()
    add_title_slide(prs)
    add_bullet_slide(prs, "Problem and Solution", [
        "Problem: prompts were scattered, hard to reuse, and error-prone.",
        "Solution: one unified platform called PrePrompt.",
        "Prompt Bank handles discovery, testing, and organization.",
        "PromptTyper handles quick insertion into typing workflows.",
    ])
    add_system_flow_slide(prs)
    add_use_case_diagram_slide(prs)
    add_results_slide(prs)
    prs.save(OUT_PATH)
    print(f"PPT_CREATED: {OUT_PATH}")


if __name__ == "__main__":
    build_presentation()
