from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT_PATH = r"C:\Users\LENOVO\Desktop\prompt bank\PrePrompt_5_Slides_TeamStyle.pptx"


def add_footer(slide, slide_no: int):
    left_text = slide.shapes.add_textbox(Inches(0.4), Inches(6.95), Inches(8.0), Inches(0.3))
    tf1 = left_text.text_frame
    tf1.text = "PrePrompt"
    p1 = tf1.paragraphs[0]
    p1.font.size = Pt(11)
    p1.font.color.rgb = RGBColor(80, 80, 80)

    right_text = slide.shapes.add_textbox(Inches(12.6), Inches(6.95), Inches(0.5), Inches(0.3))
    tf2 = right_text.text_frame
    tf2.text = str(slide_no)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(80, 80, 80)


def style_title(shape, size=38):
    tf = shape.text_frame
    for p in tf.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "PrePrompt"
    slide.placeholders[1].text = (
        "Prompt Bank and PromptTyper\n"
        "Unified Prompt Workflow System"
    )
    style_title(slide.shapes.title, 46)
    for p in slide.placeholders[1].text_frame.paragraphs:
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(40, 40, 40)

    team_box = slide.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(11.8), Inches(1.4))
    tf = team_box.text_frame
    tf.text = "PROJECT PRESENTATION"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = RGBColor(30, 30, 30)
    p2 = tf.add_paragraph()
    p2.text = "Prepared for project review"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(60, 60, 60)
    add_footer(slide, 1)


def add_bullets_slide(prs, slide_no, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    style_title(slide.shapes.title, 36)
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.font.size = Pt(23)
        p.font.color.rgb = RGBColor(30, 30, 30)
    add_footer(slide, slide_no)


def add_system_workflow_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = "SYSTEM WORKFLOW"
    style_title(slide.shapes.title, 36)

    labels = [
        "User",
        "Prompt Bank\n(Web App)",
        "Bridge API\n127.0.0.1:8765",
        "PromptTyper\nManager",
        "Shortcut Usage",
    ]
    left = Inches(0.5)
    top = Inches(2.7)
    w = Inches(2.3)
    h = Inches(1.0)
    gap = Inches(0.25)
    nodes = []

    for i, label in enumerate(labels):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + i * (w + gap), top, w, h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(236, 242, 255)
        box.line.color.rgb = RGBColor(70, 110, 180)
        box.text_frame.text = label
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for r in box.text_frame.paragraphs[0].runs:
            r.font.size = Pt(13)
            r.font.bold = True
            r.font.color.rgb = RGBColor(20, 45, 90)
        nodes.append(box)

    for i in range(len(nodes) - 1):
        x1 = nodes[i].left + nodes[i].width
        y = nodes[i].top + nodes[i].height // 2
        x2 = nodes[i + 1].left
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
        line.line.color.rgb = RGBColor(40, 40, 40)
        line.line.width = Pt(1.8)

    note = slide.shapes.add_textbox(Inches(0.9), Inches(5.1), Inches(11.6), Inches(1.2))
    tf = note.text_frame
    tf.text = "Storage: prompt_bank_data.json + local metadata (favorites, recent, usage)"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = RGBColor(35, 35, 35)
    add_footer(slide, slide_no)


def add_use_case_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = "USE CASE DIAGRAM"
    style_title(slide.shapes.title, 36)

    # System boundary
    boundary = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.3), Inches(1.5), Inches(8.4), Inches(4.8)
    )
    boundary.fill.solid()
    boundary.fill.fore_color.rgb = RGBColor(246, 249, 255)
    boundary.line.color.rgb = RGBColor(90, 90, 90)
    boundary.text_frame.text = "PrePrompt System"
    boundary.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    boundary.text_frame.paragraphs[0].font.size = Pt(13)
    boundary.text_frame.paragraphs[0].font.bold = True

    # Actors
    user = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2), Inches(1.5), Inches(0.9))
    user.text_frame.text = "User"
    user.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    user.fill.solid()
    user.fill.fore_color.rgb = RGBColor(233, 240, 255)
    user.line.color.rgb = RGBColor(70, 110, 180)

    pt = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.0), Inches(2.2), Inches(1.8), Inches(0.9))
    pt.text_frame.text = "PromptTyper"
    pt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    pt.fill.solid()
    pt.fill.fore_color.rgb = RGBColor(233, 240, 255)
    pt.line.color.rgb = RGBColor(70, 110, 180)

    use_cases = [
        ("Browse/Search Prompts", Inches(3.0), Inches(2.0)),
        ("Test Prompt", Inches(5.2), Inches(2.0)),
        ("Save to PromptTyper", Inches(7.4), Inches(2.0)),
        ("Manage Prompts", Inches(4.1), Inches(3.6)),
        ("Export/Import JSON", Inches(6.5), Inches(3.6)),
    ]
    ellipses = []
    for text, left, top in use_cases:
        e = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, Inches(2.2), Inches(1.0))
        e.fill.solid()
        e.fill.fore_color.rgb = RGBColor(255, 255, 255)
        e.line.color.rgb = RGBColor(70, 70, 70)
        e.text_frame.text = text
        e.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        e.text_frame.paragraphs[0].font.size = Pt(12)
        ellipses.append(e)

    for idx in (0, 1, 3):
        c = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            user.left + user.width,
            user.top + user.height // 2,
            ellipses[idx].left,
            ellipses[idx].top + ellipses[idx].height // 2,
        )
        c.line.color.rgb = RGBColor(90, 90, 90)
        c.line.width = Pt(1.5)

    for idx in (2, 3):
        c = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            ellipses[idx].left + ellipses[idx].width,
            ellipses[idx].top + ellipses[idx].height // 2,
            pt.left,
            pt.top + pt.height // 2,
        )
        c.line.color.rgb = RGBColor(90, 90, 90)
        c.line.width = Pt(1.5)

    add_footer(slide, slide_no)


def build():
    prs = Presentation()
    add_title_slide(prs)
    add_bullets_slide(prs, 2, "PROBLEM AND SOLUTION", [
        "Problem: prompts were scattered and difficult to reuse.",
        "Manual copy-paste increased errors and wasted time.",
        "Solution: PrePrompt centralizes prompt discovery and usage.",
        "Prompt Bank + PromptTyper create one smooth workflow.",
    ])
    add_system_workflow_slide(prs, 3)
    add_use_case_slide(prs, 4)
    add_bullets_slide(prs, 5, "RESULTS AND CONCLUSION", [
        "Prompt testing and saving now happen in one interface.",
        "Users can manage prompts without changing code.",
        "Shortcut-based usage is faster and practical for daily work.",
        "PrePrompt is demo-ready for project presentation.",
    ])
    prs.save(OUT_PATH)
    print(f"PPT_CREATED: {OUT_PATH}")


if __name__ == "__main__":
    build()
